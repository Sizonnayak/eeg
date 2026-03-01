"""
make_ppt.py  –  Generate P300 EEG Classification presentation
Usage:  python make_ppt.py
Output: P300_EEG_Classification.pptx  (in the same directory)
Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ────────────────────────────────────────────────────────────
DARK_BG   = RGBColor(0x1A, 0x1A, 0x2E)   # deep navy
MID_BG    = RGBColor(0x16, 0x21, 0x3E)   # slightly lighter navy
ACCENT    = RGBColor(0x0F, 0x3F, 0x6B)   # panel background
HIGHLIGHT = RGBColor(0x00, 0xB4, 0xD8)   # cyan-blue highlight
GREEN     = RGBColor(0x06, 0xD6, 0xA0)   # green for good numbers
ORANGE    = RGBColor(0xFF, 0xA0, 0x00)   # orange for caution
RED_COL   = RGBColor(0xEF, 0x47, 0x6F)   # red for bad
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xB0, 0xC4, 0xDE)   # light steel blue for body text

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # completely blank layout


# ── Helpers ───────────────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(BLANK)


def bg(slide, color=DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, l, t, w, h, bg_color=None, border_color=None, border_pt=0):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.line.width = Pt(border_pt)
    if border_color and border_pt > 0:
        shape.line.color.rgb = border_color
    else:
        shape.line.fill.background()
    if bg_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
    else:
        shape.fill.background()
    return shape


def txt(slide, text, l, t, w, h,
        font_size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.name  = "Calibri"
    return txb


def txt_lines(slide, lines, l, t, w, h,
              font_size=16, color=LIGHT, align=PP_ALIGN.LEFT, spacing=1.15):
    """Add a text box with multiple lines (list of (text, bold, color) tuples or plain strings)."""
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if isinstance(line, str):
            line_text, line_bold, line_color = line, False, color
        else:
            line_text = line[0]
            line_bold  = line[1] if len(line) > 1 else False
            line_color = line[2] if len(line) > 2 else color

        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        from pptx.util import Pt as _Pt
        from pptx.oxml.ns import qn
        from lxml import etree
        pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
        lnspc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPct = etree.SubElement(lnspc, qn('a:spcPct'))
        spcPct.set('val', str(int(spacing * 100000)))

        run = p.add_run()
        run.text = line_text
        run.font.size  = _Pt(font_size)
        run.font.bold  = line_bold
        run.font.color.rgb = line_color
        run.font.name  = "Calibri"
    return txb


def accent_bar(slide, color=HIGHLIGHT):
    """Thin coloured bar at the very top."""
    b = box(slide, 0, 0, 13.33, 0.06, bg_color=color)
    return b


def slide_title(slide, title, subtitle=None):
    accent_bar(slide)
    txt(slide, title, 0.4, 0.12, 12.5, 0.65,
        font_size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        txt(slide, subtitle, 0.4, 0.75, 12.5, 0.4,
            font_size=16, bold=False, color=HIGHLIGHT, align=PP_ALIGN.LEFT)


def panel(slide, l, t, w, h, title=None, title_color=HIGHLIGHT):
    box(slide, l, t, w, h, bg_color=ACCENT, border_color=HIGHLIGHT, border_pt=1)
    if title:
        txt(slide, title, l+0.1, t+0.08, w-0.2, 0.35,
            font_size=14, bold=True, color=title_color)
    return (l+0.15, t + (0.45 if title else 0.15))  # content start (l, t)


def metric_card(slide, l, t, w, h, label, value, value_color=GREEN, sub=None):
    box(slide, l, t, w, h, bg_color=RGBColor(0x0A, 0x29, 0x4A),
        border_color=HIGHLIGHT, border_pt=1)
    txt(slide, label, l+0.05, t+0.08, w-0.1, 0.3,
        font_size=11, bold=False, color=LIGHT, align=PP_ALIGN.CENTER)
    txt(slide, value, l+0.05, t+0.32, w-0.1, 0.55,
        font_size=26, bold=True, color=value_color, align=PP_ALIGN.CENTER)
    if sub:
        txt(slide, sub, l+0.05, t+0.82, w-0.1, 0.25,
            font_size=10, bold=False, color=LIGHT, align=PP_ALIGN.CENTER)


def footer(slide, text="M.Tech Project  |  P300 EEG Classification  |  2025-26"):
    box(slide, 0, 7.3, 13.33, 0.2, bg_color=RGBColor(0x0A, 0x10, 0x20))
    txt(slide, text, 0.3, 7.3, 12.7, 0.2,
        font_size=9, color=RGBColor(0x60, 0x70, 0x80), align=PP_ALIGN.LEFT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
accent_bar(sl, HIGHLIGHT)

# large title
txt(sl, "P300 EEG-Based BCI", 0.6, 1.2, 12.0, 1.1,
    font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(sl, "Cross-Subject Classification Pipeline", 0.6, 2.2, 12.0, 0.7,
    font_size=28, bold=False, color=HIGHLIGHT, align=PP_ALIGN.CENTER)

# divider
box(sl, 3.5, 3.05, 6.3, 0.04, bg_color=HIGHLIGHT)

txt(sl, "v1  →  v3  →  v4  →  Threshold Calibration", 0.6, 3.2, 12.0, 0.5,
    font_size=18, bold=False, color=LIGHT, align=PP_ALIGN.CENTER)

txt(sl, "M.Tech Project  |  Department of AI & Data Science", 0.6, 4.2, 12.0, 0.4,
    font_size=14, bold=False, color=LIGHT, align=PP_ALIGN.CENTER)
txt(sl, "2025 – 2026", 0.6, 4.6, 12.0, 0.35,
    font_size=13, bold=False, color=RGBColor(0x70, 0x80, 0x90), align=PP_ALIGN.CENTER)

footer(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — What is P300?
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
slide_title(sl, "What is P300?", "An ERP component used as a brain–computer interface signal")

# left: explanation
cl, ct = panel(sl, 0.3, 1.2, 5.8, 5.8, "The P300 Signal")
txt_lines(sl, [
    ("The P300 is a positive EEG deflection ~300 ms after a rare, attended stimulus.", False, LIGHT),
    "",
    ("In a P300 Speller:", True, HIGHLIGHT),
    ("  • A 6×6 letter matrix flashes rows and columns", False, LIGHT),
    ("  • The row/column containing the target letter evokes P300", False, LIGHT),
    ("  • Non-target flashes produce no P300", False, LIGHT),
    "",
    ("Key properties:", True, HIGHLIGHT),
    ("  • Latency: 200–500 ms post-stimulus", False, LIGHT),
    ("  • Amplitude: 5–10 µV above baseline", False, LIGHT),
    ("  • Maximal at Pz (central-parietal)", False, LIGHT),
    ("  • Highly variable across subjects", False, LIGHT),
], cl, ct, 5.4, 4.8, font_size=14)

# right: dataset info
cl2, ct2 = panel(sl, 6.4, 1.2, 6.6, 2.6, "Dataset (BCI Competition II, DS II)")
txt_lines(sl, [
    ("  8 subjects,  8 EEG channels,  250 Hz", False, LIGHT),
    ("  Epoch window:  −200 ms  to  +800 ms", False, LIGHT),
    ("  Total epochs:  33,489", False, LIGHT),
], cl2, ct2, 6.2, 1.8, font_size=14)

# class imbalance cards
metric_card(sl, 6.5, 3.95, 2.9, 1.3,
            "TARGET epochs  (y = 1)", "5,575", GREEN, "16.65% — P300-evoking")
metric_card(sl, 9.6, 3.95, 3.1, 1.3,
            "NON-TARGET  (y = 0)", "27,914", ORANGE, "83.35% — frequent flashes")

# imbalance warning
box(sl, 6.4, 5.4, 6.6, 0.95, bg_color=RGBColor(0x3A, 0x10, 0x10),
    border_color=RED_COL, border_pt=1)
txt(sl, "⚠  83:17 class imbalance — accuracy alone is misleading.\n"
        "    Use Balanced Accuracy, ROC-AUC, and PR-AUC.", 6.55, 5.48, 6.3, 0.8,
    font_size=13, color=RGBColor(0xFF, 0xCC, 0xCC))

footer(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Pipeline Overview
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
slide_title(sl, "Pipeline Overview", "Four versions built progressively to tackle cross-subject generalisation")

stages = [
    ("PREPROCESS", "0.1–30 Hz bandpass\nCommon Average Ref\nEpoch: −200 to +800 ms"),
    ("FEATURES", "Crop 200–500 ms\nBin average (5 ms)\n120-dim or Riemannian"),
    ("CLASSIFY", "rLDA  /  SVC\nStratified 5-Fold WS\nLOSO cross-subject"),
    ("CALIBRATE", "Youden's J threshold\n10% held-out val\nFix BAcc collapse"),
]
arrow_col = HIGHLIGHT
for i, (stg, desc) in enumerate(stages):
    lx = 0.4 + i * 3.22
    box(sl, lx, 1.6, 2.9, 2.8, bg_color=ACCENT, border_color=HIGHLIGHT, border_pt=1.5)
    txt(sl, stg, lx+0.1, 1.7, 2.7, 0.5, font_size=17, bold=True,
        color=HIGHLIGHT, align=PP_ALIGN.CENTER)
    box(sl, lx+1.15, 2.15, 0.6, 0.04, bg_color=HIGHLIGHT)   # mini divider
    txt(sl, desc, lx+0.1, 2.25, 2.7, 1.9, font_size=13, color=LIGHT,
        align=PP_ALIGN.CENTER)
    if i < 3:
        txt(sl, "→", lx+2.92, 2.7, 0.35, 0.5, font_size=26, bold=True,
            color=HIGHLIGHT, align=PP_ALIGN.CENTER)

# version timeline
box(sl, 0.4, 4.65, 12.5, 0.04, bg_color=RGBColor(0x30, 0x40, 0x60))
versions = [
    (0.55, "v1", "Raw bins\n120-dim rLDA"),
    (3.3,  "v3", "xDAWN sweep\nk=0,2,4,6"),
    (6.55, "v4", "Riemannian\nEA + TS + rLDA"),
    (9.8,  "v4+cal", "Youden's J\ncalibration"),
]
for vx, vname, vdesc in versions:
    box(sl, vx, 4.5, 0.08, 0.35, bg_color=HIGHLIGHT)
    txt(sl, vname, vx-0.35, 4.9, 1.2, 0.35, font_size=13, bold=True,
        color=HIGHLIGHT, align=PP_ALIGN.CENTER)
    txt(sl, vdesc, vx-0.5, 5.28, 1.5, 0.7, font_size=11, color=LIGHT,
        align=PP_ALIGN.CENTER)

footer(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Critical Bug: Label Inversion
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
slide_title(sl, "Critical Finding: Label Inversion Bug",
            "Original code assigned y=1 to NON-TARGET — silently flipped all learning")

box(sl, 0.3, 1.2, 12.7, 1.05,
    bg_color=RGBColor(0x3A, 0x08, 0x08), border_color=RED_COL, border_pt=2)
txt(sl, "flash[:, 3] encodes:   1 = NON-TARGET (frequent, 83%)     2 = TARGET (P300-evoking, 17%)",
    0.5, 1.3, 12.3, 0.45, font_size=15, bold=True, color=RGBColor(0xFF, 0xCC, 0xCC))
txt(sl, "Original code used  flash[i,3] == 1  as the positive class — the model was learning to detect NON-targets.",
    0.5, 1.65, 12.3, 0.4, font_size=13, color=RGBColor(0xFF, 0x99, 0x99))

# before / after
cl, ct = panel(sl, 0.3, 2.45, 6.0, 2.3, "BEFORE  (wrong)", RED_COL)
txt_lines(sl, [
    ("labels.append(", False, LIGHT),
    ("    1 if flash[i, 3] == 1 else 0", False, RED_COL),
    (")", False, LIGHT),
    "",
    ("→ y=1 assigned to NON-TARGET (83%)", True, RED_COL),
], cl, ct, 5.5, 1.9, font_size=14)

cl2, ct2 = panel(sl, 6.6, 2.45, 6.4, 2.3, "AFTER  (correct)", GREEN)
txt_lines(sl, [
    ("labels.append(", False, LIGHT),
    ("    1 if flash[i, 3] == 2 else 0", False, GREEN),
    (")", False, LIGHT),
    "",
    ("→ y=1 assigned to TARGET (17%)", True, GREEN),
], cl2, ct2, 5.9, 1.9, font_size=14)

txt_lines(sl, [
    ("Why it was hard to catch:", True, ORANGE),
    ("  • WS accuracy looked high (~85%) because the model just predicted the majority class", False, LIGHT),
    ("  • AUC can flip:  AUC(wrong direction) = 1 − AUC(correct direction)", False, LIGHT),
    ("  • LOSO BAcc ≈ 0.51 near chance — the real signal was invisible", False, LIGHT),
], 0.4, 4.85, 12.5, 1.4, font_size=13)

footer(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — v1 Baseline
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
slide_title(sl, "v1: Baseline — Raw Binned ERP Features",
            "Simple, interpretable; establishes the within-subject vs cross-subject gap")

# pipeline strip
box(sl, 0.3, 1.3, 12.7, 0.75, bg_color=RGBColor(0x0A, 0x20, 0x35))
steps = ["Filter\n0.1–40 Hz", "Epoch\n−200→+800ms", "Crop\n200–500ms",
         "Bin avg\n5 ms bins", "Flatten\n120-dim", "rLDA /\nSVC"]
for i, s in enumerate(steps):
    lx = 0.45 + i * 2.08
    box(sl, lx, 1.38, 1.8, 0.58, bg_color=ACCENT, border_color=HIGHLIGHT, border_pt=1)
    txt(sl, s, lx+0.05, 1.42, 1.7, 0.5, font_size=11, bold=False,
        color=LIGHT, align=PP_ALIGN.CENTER)
    if i < 5:
        txt(sl, "→", lx+1.82, 1.55, 0.28, 0.3, font_size=14, bold=True,
            color=HIGHLIGHT, align=PP_ALIGN.CENTER)

# results cards — WS
txt(sl, "Within-Subject  (5-Fold Stratified CV)", 0.4, 2.3, 6.2, 0.35,
    font_size=14, bold=True, color=HIGHLIGHT)
metric_card(sl, 0.3,  2.7, 2.9, 1.3, "Balanced Accuracy", "0.612", GREEN)
metric_card(sl, 3.35, 2.7, 2.9, 1.3, "ROC-AUC", "0.762", GREEN)
metric_card(sl, 6.4,  2.7, 2.9, 1.3, "P300 signal", "DETECTED", GREEN, "within-subject")

# results cards — LOSO
txt(sl, "Cross-Subject  (Leave-One-Subject-Out)", 0.4, 4.2, 6.2, 0.35,
    font_size=14, bold=True, color=ORANGE)
metric_card(sl, 0.3,  4.6, 2.9, 1.3, "Balanced Accuracy", "0.519", ORANGE, "≈ chance (0.500)")
metric_card(sl, 3.35, 4.6, 2.9, 1.3, "ROC-AUC", "0.636", ORANGE)
metric_card(sl, 6.4,  4.6, 2.9, 1.3, "Cross-subject gap", "LARGE", RED_COL, "BAcc drops 9.3pp")

# diagnosis
box(sl, 9.55, 2.6, 3.4, 4.3, bg_color=RGBColor(0x08, 0x18, 0x28),
    border_color=HIGHLIGHT, border_pt=1)
txt(sl, "Why LOSO fails", 9.65, 2.7, 3.2, 0.35, font_size=13, bold=True, color=HIGHLIGHT)
txt_lines(sl, [
    "• ERP latency varies ±50 ms across subjects",
    "• Amplitude scale differs 3–5×",
    "• P300 topography shifts with head shape",
    "• Threshold (0.5) from training subject",
    "  does not transfer to test subject",
    "",
    ("→ Need geometry-aware features", True, HIGHLIGHT),
    ("→ Need calibrated threshold", True, HIGHLIGHT),
], 9.65, 3.1, 3.2, 3.6, font_size=12)

footer(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — v3: xDAWN Sweep
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
slide_title(sl, "v3: xDAWN Spatial Filter Sweep  (k = 0, 2, 4, 6)",
            "Does a subject-trained spatial filter improve cross-subject generalisation?")

# table header
cols = ["k", "Feat dim", "WS BAcc", "WS AUC", "LOSO BAcc", "LOSO AUC"]
col_x = [0.3, 1.5, 3.1, 5.0, 6.9, 9.3]
col_w = [1.1, 1.5, 1.8, 1.8, 2.2, 2.0]
row_h = 0.52
header_y = 1.35

for i, (cx, cw, ch) in enumerate(zip(col_x, col_w, cols)):
    box(sl, cx, header_y, cw, row_h, bg_color=RGBColor(0x0F, 0x3F, 0x6B),
        border_color=HIGHLIGHT, border_pt=1)
    txt(sl, ch, cx+0.05, header_y+0.08, cw-0.1, row_h-0.1,
        font_size=13, bold=True, color=HIGHLIGHT, align=PP_ALIGN.CENTER)

rows = [
    ["0",  "120", "0.591", "0.742", "0.509", "0.607"],
    ["2",  " 30", "0.558", "0.719", "0.505", "0.606"],
    ["4",  " 60", "0.581", "0.732", "0.507", "0.608"],
    ["6",  " 90", "0.589", "0.740", "0.509", "0.605"],
]
for r, row in enumerate(rows):
    ry = header_y + (r+1)*row_h
    rc = ACCENT if r % 2 == 0 else RGBColor(0x12, 0x28, 0x40)
    for i, (cx, cw, val) in enumerate(zip(col_x, col_w, row)):
        box(sl, cx, ry, cw, row_h, bg_color=rc, border_color=RGBColor(0x30,0x40,0x60), border_pt=0.5)
        vc = WHITE if i < 2 else (GREEN if i in [2,3] else ORANGE)
        txt(sl, val, cx+0.05, ry+0.1, cw-0.1, row_h-0.15,
            font_size=14, bold=(i >= 4), color=vc, align=PP_ALIGN.CENTER)

# conclusion box
box(sl, 0.3, 4.5, 12.7, 1.3,
    bg_color=RGBColor(0x18, 0x10, 0x05), border_color=ORANGE, border_pt=1.5)
txt(sl, "Finding:  xDAWN does NOT improve LOSO regardless of k  —  all LOSO AUC ≈ 0.606±0.001",
    0.5, 4.58, 12.3, 0.4, font_size=15, bold=True, color=ORANGE)
txt_lines(sl, [
    ("Reason:  With only 8 EEG channels, the spatial filter learned on 7 training subjects does not capture"
     " the test subject's unique P300 topography.  The information bottleneck is topology, not feature dimension.", False, LIGHT),
], 0.5, 4.98, 12.3, 0.65, font_size=13)

txt(sl, "→ Next step: model the geometry of the covariance space (Riemannian pipeline)",
    0.4, 5.9, 12.6, 0.4, font_size=14, bold=True, color=HIGHLIGHT)

footer(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — v4: Riemannian Pipeline
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
slide_title(sl, "v4: Riemannian Covariance Pipeline",
            "XdawnCovariances → Euclidean Alignment → Tangent Space → rLDA")

# pipeline diagram
steps4 = [
    ("XdawnCov\n(nfilter=4)", "Augmented\n(2nf+1)×(2nf+1)\nSPD matrices"),
    ("Euclidean\nAlignment", "Whiten by\ntrain mean cov\nR⁻¹/² C R⁻¹/²"),
    ("Tangent\nSpace", "Map SPD →\nEuclidean\n45-dim vector"),
    ("rLDA", "Regularised\nLinear\nDiscriminant"),
]
for i, (sname, sdesc) in enumerate(steps4):
    lx = 0.35 + i * 3.22
    box(sl, lx, 1.3, 2.85, 2.1, bg_color=ACCENT, border_color=HIGHLIGHT, border_pt=1.5)
    txt(sl, sname, lx+0.1, 1.38, 2.65, 0.55, font_size=14, bold=True,
        color=HIGHLIGHT, align=PP_ALIGN.CENTER)
    box(sl, lx+0.9, 1.88, 1.05, 0.04, bg_color=RGBColor(0x30,0x50,0x70))
    txt(sl, sdesc, lx+0.1, 1.96, 2.65, 1.0, font_size=11, color=LIGHT,
        align=PP_ALIGN.CENTER)
    if i < 3:
        txt(sl, "→", lx+2.87, 1.95, 0.37, 0.5, font_size=24, bold=True,
            color=HIGHLIGHT, align=PP_ALIGN.CENTER)

# nfilter sweep table
txt(sl, "nfilter Sweep Results  (rLDA)", 0.35, 3.55, 5.0, 0.35,
    font_size=14, bold=True, color=HIGHLIGHT)

sweep_cols = ["nf", "TS dim", "WS BAcc", "WS AUC", "LOSO BAcc", "LOSO AUC", "LOSO PR-AUC"]
sc_x = [0.35, 1.35, 2.5, 4.0, 5.5, 7.15, 8.85]
sc_w = [0.9,  1.05, 1.4, 1.4, 1.55, 1.6,  1.6]
for i, (cx, cw, ch) in enumerate(zip(sc_x, sc_w, sweep_cols)):
    box(sl, cx, 3.95, cw, 0.45, bg_color=RGBColor(0x0F, 0x3F, 0x6B),
        border_color=HIGHLIGHT, border_pt=0.8)
    txt(sl, ch, cx+0.03, 4.0, cw-0.06, 0.35,
        font_size=11, bold=True, color=HIGHLIGHT, align=PP_ALIGN.CENTER)

sweep_rows = [
    ["1", " 6", "0.506", "0.582", "0.500", "0.552", "0.203"],
    ["2", "15", "0.550", "0.673", "0.504", "0.602", "0.243"],
    ["4", "45", "0.600", "0.746", "0.517", "0.648", "0.296"],
]
for r, row in enumerate(sweep_rows):
    ry = 4.4 + r * 0.45
    rc = ACCENT if r % 2 == 0 else RGBColor(0x12, 0x28, 0x40)
    best = (r == 2)
    bc = RGBColor(0x05, 0x30, 0x18) if best else rc
    for i, (cx, cw, val) in enumerate(zip(sc_x, sc_w, row)):
        box(sl, cx, ry, cw, 0.44, bg_color=bc,
            border_color=GREEN if best else RGBColor(0x30,0x40,0x60), border_pt=1 if best else 0.5)
        vc = GREEN if best else (WHITE if i < 2 else LIGHT)
        txt(sl, val, cx+0.03, ry+0.07, cw-0.06, 0.3,
            font_size=13, bold=best, color=vc, align=PP_ALIGN.CENTER)

# key insight
box(sl, 0.35, 5.85, 12.6, 1.0, bg_color=RGBColor(0x05, 0x25, 0x10),
    border_color=GREEN, border_pt=1.5)
txt(sl, "nf=4 best:  LOSO AUC 0.636 → 0.648  (+1.2pp)   PR-AUC = 0.296 = 1.78× above random (0.167)",
    0.55, 5.93, 12.2, 0.4, font_size=15, bold=True, color=GREEN)
txt(sl, "Riemannian geometry improves ranking — but LOSO BAcc still 0.517 because threshold does not transfer.",
    0.55, 6.32, 12.2, 0.35, font_size=12, color=LIGHT)

footer(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Threshold Calibration
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
slide_title(sl, "Threshold Calibration: Youden's J",
            "The AUC is good — the default threshold just doesn't transfer across subjects")

# explanation
cl, ct = panel(sl, 0.3, 1.3, 5.6, 3.5, "Why Default Threshold Fails")
txt_lines(sl, [
    "Default threshold = 0.5 works when train and",
    "test score distributions are aligned.",
    "",
    ("In LOSO:", True, HIGHLIGHT),
    "• Training subjects define one score scale",
    "• Test subject has a different scale",
    "• Threshold 0.5 is no longer the optimal cut",
    "• BAcc collapses to ~0.50 even when AUC > 0.60",
    "",
    ("Solution:  Youden's J", True, GREEN),
    "  thresh* = argmax(TPR − FPR) on ROC curve",
    "  Fitted on a 10% held-out val split of",
    "  the training subjects per LOSO fold.",
], cl, ct, 5.2, 3.0, font_size=13)

# procedure
cl2, ct2 = panel(sl, 6.2, 1.3, 6.8, 3.5, "Calibration Procedure (per LOSO fold)")
txt_lines(sl, [
    ("1. Split training subjects:", True, HIGHLIGHT),
    "   90% → fit Riemannian pipeline + rLDA",
    "   10% → find Youden's J threshold",
    "",
    ("2. Re-fit final model on 100% of training", True, HIGHLIGHT),
    "   subjects  (no data wasted)",
    "",
    ("3. Apply calibrated threshold to", True, HIGHLIGHT),
    "   held-out test subject",
    "",
    ("4. Report BAcc at calibrated threshold", True, HIGHLIGHT),
    "   (AUC unchanged — still threshold-free)",
], cl2, ct2, 6.4, 3.0, font_size=13)

# per-subject results
txt(sl, "Per-Subject Calibrated LOSO Results  (nf=4, rLDA)", 0.35, 4.95, 12.5, 0.35,
    font_size=14, bold=True, color=HIGHLIGHT)

subj_data = [
    ("S01","0.178","0.510","0.539"),
    ("S02","0.165","0.567","0.593"),
    ("S03","0.174","0.581","0.611"),
    ("S04","0.181","0.616","0.667"),
    ("S05","0.165","0.628","0.684"),
    ("S06","0.154","0.621","0.668"),
    ("S07","0.183","0.690","0.750"),
    ("S08","0.156","0.627","0.672"),
]
col_heads = ["Subj","thresh","BAcc","AUC"]
cw2 = [0.85, 0.9, 0.82, 0.82]
cx2 = [0.3 + sum(cw2[:i]) + i*0.18 for i in range(4)]

for i, (cx, cw, ch) in enumerate(zip(cx2, cw2, col_heads)):
    box(sl, cx, 5.35, cw, 0.38, bg_color=RGBColor(0x0F,0x3F,0x6B),
        border_color=HIGHLIGHT, border_pt=0.8)
    txt(sl, ch, cx+0.03, 5.38, cw-0.06, 0.3,
        font_size=11, bold=True, color=HIGHLIGHT, align=PP_ALIGN.CENTER)

for r, (subj, thresh, bacc, auc) in enumerate(subj_data):
    ry = 5.73 + r * 0.195
    rc = ACCENT if r % 2 == 0 else RGBColor(0x12, 0x28, 0x40)
    bv = float(bacc)
    bc_val = GREEN if bv >= 0.60 else (ORANGE if bv >= 0.55 else LIGHT)
    for i, (cx, cw, val) in enumerate(zip(cx2, cw2, [subj,thresh,bacc,auc])):
        box(sl, cx, ry, cw, 0.19, bg_color=rc,
            border_color=RGBColor(0x30,0x40,0x60), border_pt=0.3)
        vc = bc_val if i == 2 else LIGHT
        txt(sl, val, cx+0.02, ry+0.02, cw-0.04, 0.15,
            font_size=10, bold=(i==2), color=vc, align=PP_ALIGN.CENTER)

# final numbers
metric_card(sl, 5.2, 4.95, 2.5, 1.1, "Calibrated BAcc", "0.605", GREEN, "vs 0.517 uncalib")
metric_card(sl, 7.85, 4.95, 2.5, 1.1, "ROC-AUC", "0.648", GREEN, "unchanged")
metric_card(sl, 10.5, 4.95, 2.55, 1.1, "PR-AUC", "0.296", GREEN, "1.78× random")

footer(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Results Summary
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
slide_title(sl, "Results Summary — All Versions",
            "Progressive improvement from raw features to calibrated Riemannian pipeline")

# full comparison table
cols9 = ["Version", "Method", "Eval", "BAcc", "ROC-AUC", "PR-AUC", "Notes"]
cw9   = [1.2,        3.3,     0.8,    0.9,    1.0,       0.95,    3.55]
cx9   = [0.3 + sum(cw9[:i]) + i*0.06 for i in range(7)]

for i, (cx, cw, ch) in enumerate(zip(cx9, cw9, cols9)):
    box(sl, cx, 1.35, cw, 0.45, bg_color=RGBColor(0x0F,0x3F,0x6B),
        border_color=HIGHLIGHT, border_pt=1)
    txt(sl, ch, cx+0.05, 1.38, cw-0.1, 0.35,
        font_size=12, bold=True, color=HIGHLIGHT, align=PP_ALIGN.CENTER)

table_rows = [
    ("v1", "rLDA 120d raw",       "WS",   "0.612", "0.762", "—",     "Strong within-subject",       GREEN,  WHITE),
    ("v1", "rLDA 120d raw",       "LOSO", "0.519", "0.636", "—",     "Near-chance BAcc",             ORANGE, ORANGE),
    ("v3 k=0","rLDA 120d CAR",    "LOSO", "0.509", "0.607", "—",     "CAR/filter gives no LOSO gain",ORANGE, LIGHT),
    ("v3 k=6","rLDA+xDAWN(6)",    "LOSO", "0.509", "0.605", "—",     "xDAWN does not help LOSO",    ORANGE, LIGHT),
    ("v4 nf=4","Riem EA+TS rLDA", "WS",   "0.600", "0.746", "0.443", "Best WS after v1",             GREEN,  WHITE),
    ("v4 nf=4","Riem EA+TS rLDA", "LOSO", "0.517", "0.648", "0.296", "Best AUC, threshold issue",    ORANGE, ORANGE),
    ("v4+cal","Riem+Youden's J",  "LOSO", "0.605", "0.648", "0.296", "BEST cross-subject system ★",  GREEN,  GREEN),
]

for r, (ver, mth, ev, bac, auc, pr, note, bac_col, note_col) in enumerate(table_rows):
    ry = 1.8 + r * 0.52
    best = (r == 6)
    rc = RGBColor(0x05, 0x30, 0x18) if best else (ACCENT if r % 2 == 0 else RGBColor(0x12,0x28,0x40))
    bc = GREEN if best else RGBColor(0x30,0x40,0x60)
    vals = [ver, mth, ev, bac, auc, pr, note]
    for i, (cx, cw, val) in enumerate(zip(cx9, cw9, vals)):
        box(sl, cx, ry, cw, 0.5, bg_color=rc, border_color=bc, border_pt=1 if best else 0.4)
        vc = bac_col if i == 3 else (note_col if i == 6 else (WHITE if i < 3 else LIGHT))
        txt(sl, val, cx+0.04, ry+0.1, cw-0.08, 0.3,
            font_size=11, bold=best, color=vc, align=PP_ALIGN.LEFT if i in [1,6] else PP_ALIGN.CENTER)

# random baseline
txt(sl, "Random baselines:  BAcc=0.500   AUC=0.500   PR-AUC=0.167 (= prevalence 16.65%)",
    0.35, 5.55, 12.6, 0.35, font_size=12, color=RGBColor(0x80,0x90,0xA0))

footer(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Metrics Explained
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
slide_title(sl, "Why These Metrics?",
            "Accuracy is misleading under 83:17 imbalance — three complementary metrics used")

metrics = [
    ("Balanced Accuracy\n(BAcc)", "= (TPR + TNR) / 2\n\nEqual weight to both classes.\nRandom baseline = 0.500.\nOur best = 0.605.", "Measures threshold\nperformance fairly.", GREEN),
    ("ROC-AUC", "Area under ROC curve.\n\nThreshold-free ranking metric.\nRandom baseline = 0.500.\nOur best LOSO = 0.648.", "Does the model rank\ntargets above non-targets?", HIGHLIGHT),
    ("PR-AUC\n(Avg Precision)", "Area under Precision-Recall.\n\nSensitive to minority class.\nRandom baseline = prevalence\n= 0.1665.\nOur best = 0.296 (1.78×).", "Most honest metric\nunder severe imbalance.", ORANGE),
]

for i, (name, formula, insight, col) in enumerate(metrics):
    lx = 0.35 + i * 4.33
    box(sl, lx, 1.35, 4.05, 5.5, bg_color=ACCENT, border_color=col, border_pt=2)
    txt(sl, name, lx+0.1, 1.45, 3.85, 0.65, font_size=16, bold=True,
        color=col, align=PP_ALIGN.CENTER)
    box(sl, lx+0.5, 2.08, 3.05, 0.04, bg_color=col)
    txt(sl, formula, lx+0.15, 2.18, 3.75, 2.2, font_size=13, color=LIGHT)
    box(sl, lx+0.1, 4.4, 3.85, 0.04, bg_color=RGBColor(0x30,0x40,0x60))
    txt(sl, insight, lx+0.15, 4.52, 3.75, 0.9, font_size=13, bold=True,
        color=col, align=PP_ALIGN.CENTER)

txt(sl, "All three metrics are reported — they measure different aspects of the same classifier.",
    0.35, 7.0, 12.6, 0.35, font_size=13, bold=True, color=LIGHT, align=PP_ALIGN.CENTER)

footer(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Conclusions & Future Work
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
slide_title(sl, "Conclusions & Future Work")

# conclusions
cl, ct = panel(sl, 0.3, 1.3, 6.1, 5.5, "Conclusions")
txt_lines(sl, [
    ("1.  Label sanity check is critical", True, ORANGE),
    "    flash[:,3]==2 is TARGET; original code had it inverted.",
    "",
    ("2.  Within-subject P300 is detectable", True, GREEN),
    "    v1 rLDA:  WS BAcc=0.612,  AUC=0.762",
    "",
    ("3.  xDAWN does not help cross-subject", True, ORANGE),
    "    8-channel spatial filter doesn't transfer (v3 sweep).",
    "",
    ("4.  Riemannian geometry improves ranking", True, GREEN),
    "    LOSO AUC: 0.636 → 0.648 with EA + TangentSpace.",
    "",
    ("5.  Threshold calibration recovers BAcc", True, GREEN),
    "    Youden's J:  BAcc 0.517 → 0.605  (+8.8 pp)",
    "",
    ("Best system:  v4 nf=4 + Youden's J", True, HIGHLIGHT),
    "  LOSO  BAcc=0.605  AUC=0.648  PR=0.296 (1.78× chance)",
], cl, ct, 5.7, 5.0, font_size=13)

# future work
cl2, ct2 = panel(sl, 6.65, 1.3, 6.3, 5.5, "Future Work")
txt_lines(sl, [
    ("Short term:", True, HIGHLIGHT),
    "  • Sweep nfilter to 6, 8",
    "  • Try MDM (Minimum Distance to Mean)",
    "    as Riemannian classifier",
    "  • Platt scaling as alternative calibration",
    "",
    ("Medium term:", True, HIGHLIGHT),
    "  • Riemannian alignment at test time",
    "    (online EA with sliding window)",
    "  • Subject-independent model fine-tuning",
    "    with minimal labelled test data (10 trials)",
    "",
    ("Long term:", True, HIGHLIGHT),
    "  • Deep learning: EEGNet, ShallowConvNet",
    "  • Transfer learning across datasets",
    "  • Real-time P300 speller evaluation",
    "  • More subjects / higher-density EEG",
], cl2, ct2, 5.9, 4.9, font_size=13)

footer(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Dataset Details
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
slide_title(sl, "Dataset: P300 Speller — 8 Subjects",
            "MATLAB .mat files, 250 Hz, 8 EEG channels, ~24 min recording per subject")

# Left: dataset metadata table
txt(sl, "Per-Subject Statistics", 0.3, 1.3, 6.0, 0.35,
    font_size=14, bold=True, color=HIGHLIGHT)

ds_heads = ["Subject", "Epochs", "Target", "Non-target", "Bad Chs", "Threshold"]
ds_cw    = [1.0, 0.9, 0.85, 1.1, 1.05, 1.1]
ds_cx    = [0.3 + sum(ds_cw[:i]) + i*0.03 for i in range(6)]

for i, (cx, cw, ch) in enumerate(zip(ds_cx, ds_cw, ds_heads)):
    box(sl, cx, 1.68, cw, 0.38, bg_color=RGBColor(0x0F,0x3F,0x6B),
        border_color=HIGHLIGHT, border_pt=0.8)
    txt(sl, ch, cx+0.03, 1.7, cw-0.06, 0.3,
        font_size=11, bold=True, color=HIGHLIGHT, align=PP_ALIGN.CENTER)

ds_rows = [
    ("S01", "4200", "700", "3500",  "—",         "150 µV"),
    ("S02", "4198", "700", "3498",  "—",         "200 µV"),
    ("S03", "4197", "700", "3497",  "P4, PO8",   "100 µV"),
    ("S04", "4191", "697", "3494",  "PO8",       "100 µV"),
    ("S05", "4112", "680", "3432",  "Oz",        "100 µV"),
    ("S06", "4197", "700", "3497",  "—",         "150 µV"),
    ("S07", "4198", "698", "3500",  "—",         "150 µV"),
    ("S08", "4196", "700", "3496",  "—",         "150 µV"),
]
for r, row in enumerate(ds_rows):
    ry = 2.06 + r * 0.42
    rc = ACCENT if r % 2 == 0 else RGBColor(0x12, 0x28, 0x40)
    bad_flag = row[4] != "—"
    for i, (cx, cw, val) in enumerate(zip(ds_cx, ds_cw, row)):
        box(sl, cx, ry, cw, 0.4, bg_color=rc,
            border_color=RGBColor(0x30,0x40,0x60), border_pt=0.4)
        vc = ORANGE if (i == 4 and bad_flag) else LIGHT
        txt(sl, val, cx+0.03, ry+0.07, cw-0.06, 0.26,
            font_size=11, color=vc, align=PP_ALIGN.CENTER)

# total row
ry = 2.06 + 8 * 0.42
box(sl, ds_cx[0], ry, sum(ds_cw) + 5*0.03, 0.4,
    bg_color=RGBColor(0x0F,0x3F,0x6B), border_color=HIGHLIGHT, border_pt=1)
txt(sl, "TOTAL   33,489 epochs     5,575 target (16.65%)     27,914 non-target (83.35%)",
    ds_cx[0]+0.1, ry+0.08, 5.8, 0.26,
    font_size=11, bold=True, color=GREEN, align=PP_ALIGN.LEFT)

# Right: raw file structure
cl2, ct2 = panel(sl, 6.45, 1.3, 6.5, 3.55, "Raw .mat File Structure")
txt_lines(sl, [
    ("scipy.io.loadmat('P300S01.mat')", False, HIGHLIGHT),
    "",
    ("  mat['data']['X']     → (358372, 8)", False, GREEN),
    "  Continuous EEG  [samples × channels]",
    "  ~24 min @ 250 Hz",
    "",
    ("  mat['data']['flash'] → (21495, 4)", False, GREEN),
    "  Col 0: onset sample index",
    "  Col 1: row/col ID (12 rows+cols in 6×6 grid)",
    ("  Col 3: 1=NON-TARGET  2=TARGET  ← key!", True, ORANGE),
    "",
    ("  mat['Fs']        = 250 Hz", False, LIGHT),
    ("  mat['Trials']    = 35 words", False, LIGHT),
    ("  mat['Word']      = 'TOKENMIRARJUJU...'", False, LIGHT),
    ("  mat['channelNames'] = [FZ,CZ,P3,PZ,P4,PO7,PO8,OZ]", False, LIGHT),
], cl2, ct2, 6.2, 3.0, font_size=12)

# Channel map note
box(sl, 6.45, 5.0, 6.5, 1.15, bg_color=RGBColor(0x05,0x25,0x10),
    border_color=GREEN, border_pt=1)
txt(sl, "EEG Channel Layout  (8 channels, 10-20 system)", 6.6, 5.07, 6.2, 0.3,
    font_size=12, bold=True, color=GREEN)
txt_lines(sl, [
    "Frontal: Fz, Cz       Parietal: P3, Pz, P4",
    "Parieto-occipital: PO7, PO8       Occipital: Oz",
    "→ Coverage chosen for P300 (max amplitude at Pz/Cz)",
], 6.6, 5.38, 6.2, 0.7, font_size=11)

footer(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — End-to-End Pipeline: Data Loading → Preprocessing → Features
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
slide_title(sl, "End-to-End Flow: Loading → Preprocessing → Features",
            "Every transformation step with exact shapes and parameters")

# ── STEP 1: Data Loading ──────────────────────────────────────────────────
box(sl, 0.3, 1.25, 12.7, 0.3, bg_color=RGBColor(0x0F,0x3F,0x6B))
txt(sl, "STEP 1 — Data Loading  (scipy.io.loadmat)", 0.45, 1.28, 7.0, 0.24,
    font_size=12, bold=True, color=HIGHLIGHT)
txt(sl, "P300S01.mat … P300S08.mat", 8.5, 1.28, 4.3, 0.24,
    font_size=11, color=LIGHT, align=PP_ALIGN.RIGHT)

txt_lines(sl, [
    ("mat = scipy.io.loadmat('P300S0X.mat')", False, GREEN),
    "X_raw = mat['data']['X']       # (358372, 8)  continuous EEG  µV",
    "flash = mat['data']['flash']   # (21495, 4)   stimulus event table",
    "onset = flash[:, 0]            # sample index of each flash",
    "label = flash[:, 3]            # 1=non-target  2=target",
], 0.35, 1.6, 6.5, 1.1, font_size=11)

txt_lines(sl, [
    ("Output per subject:", True, HIGHLIGHT),
    "  X_raw : (N_samples, 8)  raw continuous signal",
    "  flash : (N_flashes, 4)  event markers",
    "  ~21,500 flash events per subject",
    "  ~358,000 raw samples (~24 min)",
], 7.1, 1.6, 5.8, 1.1, font_size=11)

# ── STEP 2: Preprocessing ─────────────────────────────────────────────────
box(sl, 0.3, 2.78, 12.7, 0.3, bg_color=RGBColor(0x0F,0x3F,0x6B))
txt(sl, "STEP 2 — Preprocessing  (MNE-Python)", 0.45, 2.81, 7.0, 0.24,
    font_size=12, bold=True, color=HIGHLIGHT)
txt(sl, "p300_preprocess_v2.py", 8.5, 2.81, 4.3, 0.24,
    font_size=11, color=LIGHT, align=PP_ALIGN.RIGHT)

preproc_steps = [
    ("2a", "Bad channel interpolation", "S03→P4,PO8   S04→PO8   S05→Oz\nSpherical spline interpolation from neighbours", ORANGE),
    ("2b", "Bandpass filter",           "0.1 – 30 Hz  (FIR, Hamming window)\nRemoves DC drift + high-freq noise + muscle artefacts", HIGHLIGHT),
    ("2c", "Common Average Reference",  "Each sample -= mean across all 8 channels\nReduces global noise / electrode drift shared across channels", HIGHLIGHT),
    ("2d", "Epoching",                  "Onset ± window: −200 ms to +800 ms = 250 samples per epoch\n(X_raw[onset−50 : onset+200, :])", GREEN),
    ("2e", "Baseline correction",       "Subtract mean of pre-stimulus window [−200 to 0 ms] per channel\nRemoves DC offset within each epoch", GREEN),
    ("2f", "Artifact rejection",        "Reject epoch if any sample > threshold (100–200 µV, subject-specific)\nS01 keeps ~4200, S05 most rejected (threshold 100 µV)", RED_COL),
]
step_w = 2.06
for i, (code, name, desc, col) in enumerate(preproc_steps):
    lx = 0.3 + i * (step_w + 0.04)
    box(sl, lx, 3.15, step_w, 2.1, bg_color=ACCENT, border_color=col, border_pt=1.2)
    txt(sl, code, lx+0.08, 3.2, 0.4, 0.28, font_size=11, bold=True, color=col)
    txt(sl, name, lx+0.08, 3.48, step_w-0.16, 0.35, font_size=11, bold=True, color=WHITE)
    txt(sl, desc, lx+0.08, 3.83, step_w-0.16, 1.3, font_size=10, color=LIGHT)

# shape annotation
txt(sl, "Shape after epoching:  (n_epochs, 8, 250)   →   33,489 epochs total across 8 subjects",
    0.35, 5.32, 12.5, 0.3, font_size=12, bold=True, color=GREEN)

# ── STEP 3: Feature Extraction ────────────────────────────────────────────
box(sl, 0.3, 5.68, 12.7, 0.3, bg_color=RGBColor(0x0F,0x3F,0x6B))
txt(sl, "STEP 3 — Feature Extraction  (v1/v3 path)", 0.45, 5.71, 7.0, 0.24,
    font_size=12, bold=True, color=HIGHLIGHT)
txt(sl, "extract_p300_features_v2.py", 8.5, 5.71, 4.3, 0.24,
    font_size=11, color=LIGHT, align=PP_ALIGN.RIGHT)

feat_steps = [
    ("Crop P300\nwindow", "(n,8,250)\n→(n,8,75)", "200–500ms\nsamples 100:175"),
    ("Bin\naveraging", "(n,8,75)\n→(n,8,15)", "15 bins\n× 5 samples each"),
    ("Flatten", "(n,8,15)\n→(n,120)", "8 channels\n× 15 bins"),
    ("Label\nassignment", "flash[:,3]==2\n→ y=1 target", "flash[:,3]==1\n→ y=0 non-tgt"),
]
for i, (name, shape, note) in enumerate(feat_steps):
    lx = 0.35 + i * 3.2
    box(sl, lx, 6.05, 2.9, 1.18, bg_color=ACCENT, border_color=HIGHLIGHT, border_pt=1)
    txt(sl, name, lx+0.1, 6.1, 2.7, 0.4, font_size=11, bold=True,
        color=HIGHLIGHT, align=PP_ALIGN.CENTER)
    txt(sl, shape, lx+0.1, 6.48, 2.7, 0.35, font_size=11, bold=True,
        color=GREEN, align=PP_ALIGN.CENTER)
    txt(sl, note, lx+0.1, 6.82, 2.7, 0.3, font_size=10,
        color=LIGHT, align=PP_ALIGN.CENTER)
    if i < 3:
        txt(sl, "→", lx+2.92, 6.5, 0.3, 0.3, font_size=20, bold=True,
            color=HIGHLIGHT, align=PP_ALIGN.CENTER)

footer(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — End-to-End Flow: Classification → Evaluation
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
slide_title(sl, "End-to-End Flow: Classification → Evaluation",
            "Two evaluation protocols: within-subject 5-fold CV and cross-subject LOSO")

# ── STEP 4: v4 Riemannian Classification ──────────────────────────────────
box(sl, 0.3, 1.25, 12.7, 0.3, bg_color=RGBColor(0x0F,0x3F,0x6B))
txt(sl, "STEP 4 — v4 Riemannian Classification  (inside each CV fold)", 0.45, 1.28, 9.0, 0.24,
    font_size=12, bold=True, color=HIGHLIGHT)
txt(sl, "lda_svc_classifier_v4.py", 9.8, 1.28, 3.0, 0.24,
    font_size=11, color=LIGHT, align=PP_ALIGN.RIGHT)

v4_steps = [
    ("XdawnCov\n(pyriemann)", "(n,8,250)\n→(n,9,9)", "nfilter=4\nFitted on train\nonly (no leakage)"),
    ("Euclidean\nAlignment", "(n,9,9)\n→(n,9,9)", "Whiten by\ntrain mean cov\nR⁻¹/² C R⁻¹/²"),
    ("Tangent\nSpace", "(n,9,9)\n→(n,45)", "Riemannian log\nat mean SPD\nn(n+1)/2 = 45"),
    ("Standard\nScaler", "(n,45)\n→(n,45)", "Zero mean\nUnit variance\nFitted on train"),
    ("rLDA /\nSVC", "(n,45)\n→score", "Shrinkage LDA\nor linear SVC\npredict_proba"),
]
for i, (name, shape, note) in enumerate(v4_steps):
    lx = 0.35 + i * 2.58
    box(sl, lx, 1.62, 2.32, 1.85, bg_color=ACCENT, border_color=HIGHLIGHT, border_pt=1)
    txt(sl, name, lx+0.08, 1.68, 2.16, 0.42, font_size=11, bold=True,
        color=HIGHLIGHT, align=PP_ALIGN.CENTER)
    txt(sl, shape, lx+0.08, 2.08, 2.16, 0.35, font_size=11, bold=True,
        color=GREEN, align=PP_ALIGN.CENTER)
    txt(sl, note,  lx+0.08, 2.42, 2.16, 0.9,  font_size=10,
        color=LIGHT, align=PP_ALIGN.CENTER)
    if i < 4:
        txt(sl, "→", lx+2.34, 2.35, 0.26, 0.3, font_size=18, bold=True,
            color=HIGHLIGHT, align=PP_ALIGN.CENTER)

# ── STEP 5: Evaluation Protocols ──────────────────────────────────────────
box(sl, 0.3, 3.6, 12.7, 0.3, bg_color=RGBColor(0x0F,0x3F,0x6B))
txt(sl, "STEP 5 — Evaluation Protocols", 0.45, 3.63, 7.0, 0.24,
    font_size=12, bold=True, color=HIGHLIGHT)

# WS panel
cl, ct = panel(sl, 0.3, 3.97, 5.9, 2.15, "Within-Subject CV  (WS)")
txt_lines(sl, [
    ("StratifiedKFold(n_splits=5)", False, GREEN),
    "",
    "  • Split each subject's epochs 80/20",
    "  • Train + test on SAME subject",
    "  • Measures: upper bound performance",
    "  • Stratified → preserves 16.7% ratio",
    "  • Metrics averaged across 5 folds",
], cl, ct, 5.5, 1.75, font_size=12)

# LOSO panel
cl2, ct2 = panel(sl, 6.5, 3.97, 6.5, 2.15, "Leave-One-Subject-Out  (LOSO)")
txt_lines(sl, [
    ("LeaveOneGroupOut()  — 8 folds", False, GREEN),
    "",
    "  • Train on 7 subjects  →  Test on 1",
    "  • Each subject is test subject once",
    "  • Measures: cross-subject generalisation",
    "  • Hardest evaluation — no subject overlap",
    "  • Threshold calibrated via Youden's J (v4)",
], cl2, ct2, 6.2, 1.75, font_size=12)

# ── STEP 6: Metrics ────────────────────────────────────────────────────────
box(sl, 0.3, 6.2, 12.7, 0.3, bg_color=RGBColor(0x0F,0x3F,0x6B))
txt(sl, "STEP 6 — Metrics Reported", 0.45, 6.23, 7.0, 0.24,
    font_size=12, bold=True, color=HIGHLIGHT)

metrics_info = [
    ("BAcc", "(TPR+TNR)/2", "Fair under imbalance\nRandom = 0.500", GREEN),
    ("ROC-AUC", "roc_auc_score\n(y, proba[:,1])", "Threshold-free\nRandom = 0.500", HIGHLIGHT),
    ("PR-AUC", "average_precision\n_score", "Minority-class sensitive\nRandom = 0.167", ORANGE),
    ("F1", "2·P·R/(P+R)\non target class", "Calibrated result only\nLow under imbalance", LIGHT),
]
for i, (name, formula, note, col) in enumerate(metrics_info):
    lx = 0.35 + i * 3.22
    box(sl, lx, 6.57, 2.98, 0.75, bg_color=ACCENT, border_color=col, border_pt=1)
    txt(sl, name,    lx+0.08, 6.6,  2.82, 0.25, font_size=12, bold=True,  color=col,   align=PP_ALIGN.CENTER)
    txt(sl, formula, lx+0.08, 6.84, 2.82, 0.22, font_size=9,  bold=False, color=GREEN, align=PP_ALIGN.CENTER)
    txt(sl, note,    lx+0.08, 7.03, 2.82, 0.2,  font_size=9,  bold=False, color=LIGHT, align=PP_ALIGN.CENTER)

footer(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Thank You
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
accent_bar(sl, HIGHLIGHT)
box(sl, 0, 7.3, 13.33, 0.2, bg_color=RGBColor(0x0A, 0x10, 0x20))

txt(sl, "Thank You", 0.6, 1.5, 12.0, 1.2,
    font_size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
box(sl, 3.5, 2.8, 6.3, 0.05, bg_color=HIGHLIGHT)

txt(sl, "Questions?", 0.6, 3.0, 12.0, 0.65,
    font_size=28, color=HIGHLIGHT, align=PP_ALIGN.CENTER)

# summary strip
box(sl, 1.0, 4.0, 11.33, 1.8, bg_color=RGBColor(0x0A, 0x22, 0x38),
    border_color=HIGHLIGHT, border_pt=1)
summary_metrics = [
    ("Best WS", "BAcc 0.612\nAUC 0.762", GREEN),
    ("Best LOSO AUC", "0.648\n(nf=4 Riem)", HIGHLIGHT),
    ("Calibrated LOSO", "BAcc 0.605\nAUC 0.648", GREEN),
    ("PR-AUC lift", "0.296\n1.78× chance", ORANGE),
]
for i, (label, val, col) in enumerate(summary_metrics):
    lx = 1.2 + i * 2.75
    txt(sl, label, lx, 4.1, 2.5, 0.35, font_size=12, color=LIGHT, align=PP_ALIGN.CENTER)
    txt(sl, val, lx, 4.45, 2.5, 0.9, font_size=18, bold=True, color=col, align=PP_ALIGN.CENTER)

txt(sl, "M.Tech Project  |  P300 EEG-Based BCI  |  2025–26",
    0.6, 6.0, 12.0, 0.4, font_size=14, color=LIGHT, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────────
OUT = "/Users/siznayak/Documents/others/MTech/EEG_Classification/P300_EEG_Classification.pptx"
prs.save(OUT)
print(f"Saved → {OUT}")
print("Slides: 15")
print("  1.  Title")
print("  2.  What is P300?")
print("  3.  Pipeline Overview")
print("  4.  Critical Finding: Label Inversion Bug")
print("  5.  v1 Baseline")
print("  6.  v3 xDAWN Sweep")
print("  7.  v4 Riemannian Pipeline")
print("  8.  Threshold Calibration")
print("  9.  Results Summary")
print(" 10.  Metrics Explained")
print(" 11.  Conclusions & Future Work")
print(" 12.  Dataset Details (per-subject stats + raw .mat structure)")
print(" 13.  End-to-End Flow: Loading → Preprocessing → Features")
print(" 14.  End-to-End Flow: Classification → Evaluation")
print(" 15.  Thank You")
